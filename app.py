import streamlit as st
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_SHAPE
from pptx.util import Inches
from pptx.dml.color import RGBColor
import io
import time
import re
from PIL import Image
from google import genai
from google.genai import types
from lxml import etree 

# --- Global State Initialization ---
if 'caption_cache' not in st.session_state:
    st.session_state.caption_cache = {}
if 'last_api_call' not in st.session_state:
    st.session_state.last_api_call = 0.0

# --- Advanced Text Extraction ---
def get_shape_text(shape):
    text_content = []
    
    if getattr(shape, "has_text_frame", False):
        for paragraph in shape.text_frame.paragraphs:
            text_content.append(paragraph.text)
            
    elif getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
        try:
            for subshape in shape.shapes:
                text_content.append(get_shape_text(subshape))
        except AttributeError:
            pass
            
    elif getattr(shape, "shape_type", None) in [19, 24] or hasattr(shape._element, 'nvGraphicFramePr'):
        try:
            for text_node in shape._element.xpath('.//a:t'):
                if text_node.text and text_node.text.strip():
                    text_content.append(text_node.text.strip())
                    
            for rel in shape.part.rels.values():
                if "diagramData" in rel.reltype:
                    xml_str = rel.target_part.blob.decode('utf-8', errors='ignore')
                    hidden_texts = re.findall(r'<a:t[^>]*>(.*?)</a:t>', xml_str)
                    text_content.extend([t.strip() for t in hidden_texts if t.strip()])
        except Exception:
            pass
        
    return " ".join(text_content).strip()

def get_slide_text(slide):
    text_runs = []
    for shape in slide.shapes:
        text_runs.append(get_shape_text(shape))
    return " ".join(text_runs).strip()

# --- Native ADA XML Injection Helpers ---
def set_alt_text(shape, alt_text):
    try:
        for prop in ['nvPicPr', 'nvSpPr', 'nvGraphicFramePr', 'nvGrpSpPr']:
            if hasattr(shape._element, prop):
                cNvPr = getattr(shape._element, prop).cNvPr
                cNvPr.set('descr', alt_text)
                return
    except Exception:
        pass

def mark_as_decorative(shape):
    try:
        cNvPr = None
        for prop in ['nvPicPr', 'nvSpPr', 'nvGraphicFramePr', 'nvGrpSpPr']:
            if hasattr(shape._element, prop):
                cNvPr = getattr(shape._element, prop).cNvPr
                break
        if cNvPr is None: return

        ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        extLst = cNvPr.find(f'.//{{{ns_a}}}extLst')
        if extLst is None:
            extLst = etree.SubElement(cNvPr, f'{{{ns_a}}}extLst')

        dec_uri = "{C183D7F6-B498-43B3-948B-1728B52AA6E4}"
        dec_ext = extLst.find(f'.//{{{ns_a}}}ext[@uri="{dec_uri}"]')
        if dec_ext is None:
            ext = etree.SubElement(extLst, f'{{{ns_a}}}ext', uri=dec_uri)
            etree.SubElement(ext, '{http://schemas.microsoft.com/office/drawing/2017/decorative}decorative', val="1")
            
        cNvPr.set('descr', 'DECORATIVE')
    except Exception:
        pass

def mute_smartart_children(shape):
    try:
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
            for subshape in shape.shapes:
                mark_as_decorative(subshape)
                mute_smartart_children(subshape) 
                
        elif getattr(shape, "shape_type", None) in [19, 24] or hasattr(shape._element, 'nvGraphicFramePr'):
            graphicData = shape._element.xpath('.//p:graphicData')
            if graphicData:
                child_nvPrs = graphicData[0].xpath('.//p:cNvPr', namespaces={'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'})
                ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
                
                for cNvPr in child_nvPrs:
                    cNvPr.set('descr', 'DECORATIVE') 
                    extLst = cNvPr.find(f'.//{{{ns_a}}}extLst')
                    if extLst is None:
                        extLst = etree.SubElement(cNvPr, f'{{{ns_a}}}extLst')
                    
                    dec_uri = "{C183D7F6-B498-43B3-948B-1728B52AA6E4}"
                    if extLst.find(f'.//{{{ns_a}}}ext[@uri="{dec_uri}"]') is None:
                        ext = etree.SubElement(extLst, f'{{{ns_a}}}ext', uri=dec_uri)
                        etree.SubElement(ext, '{http://schemas.microsoft.com/office/drawing/2017/decorative}decorative', val="1")
    except Exception:
        pass

# --- The Safe Ghost Overlay Logic ---
def safe_get_coords(shape):
    try:
        return shape.left, shape.top, shape.width, shape.height
    except Exception:
        try:
            xfrm = shape._element.xpath('.//p:xfrm | .//a:xfrm')[0]
            off = xfrm.xpath('.//a:off')[0]
            ext = xfrm.xpath('.//a:ext')[0]
            return int(off.get('x')), int(off.get('y')), int(ext.get('cx')), int(ext.get('cy'))
        except Exception:
            return Inches(1), Inches(1), Inches(8), Inches(4)

def create_ghost_overlay(slide, shape, caption):
    try:
        left, top, width, height = safe_get_coords(shape)
        
        overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        
        # 1. UNIQUE NAMING (Fixes PDF Merger Bug)
        overlay.name = f"ADA_Ghost_Overlay_{shape.shape_id}"
        
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = RGBColor(255, 255, 255)
        
        # 2. BULLETPROOF TRANSPARENCY (Fixes Solid White Bug)
        try:
            spPr = overlay._element.spPr
            ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
            srgbClr = spPr.xpath('.//a:solidFill/a:srgbClr', namespaces=ns)
            if srgbClr:
                alpha = etree.SubElement(srgbClr[0], '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
                alpha.set('val', '1000') # 1000 = 1% opacity
        except Exception as e:
            pass # Failsafe
            
        overlay.line.fill.background()
        set_alt_text(overlay, caption)
    except Exception:
        pass

def force_textbox_to_title(txBox):
    try:
        nvSpPr = txBox._element.nvSpPr
        cNvSpPr = nvSpPr.cNvSpPr
        if 'txBox' in cNvSpPr.attrib:
            del cNvSpPr.attrib['txBox']
            
        nvSpPr.cNvPr.set('name', 'Title 1')
        
        nvPr = nvSpPr.nvPr
        for child in nvPr.getchildren():
            nvPr.remove(child)
        ns_p = "http://schemas.openxmlformats.org/presentationml/2006/main"
        etree.SubElement(nvPr, f"{{{ns_p}}}ph", type="title")
    except Exception:
        pass

def fix_reading_order(slide):
    shapes = list(slide.shapes)
    sortable_shapes = []
    for shape in shapes:
        try:
            if hasattr(shape, "top") and hasattr(shape, "left"):
                if shape.top is not None and shape.left is not None:
                    sortable_shapes.append(shape)
        except Exception:
            pass
            
    sortable_shapes.sort(key=lambda s: (round(s.top / 100000) * 100000, s.left))
    
    if sortable_shapes:
        parent = sortable_shapes[0]._element.getparent()
        for shape in sortable_shapes:
            parent.remove(shape._element)
            parent.append(shape._element)

# --- AI Generation ---
def generate_caption(client, image_bytes, prev_text, curr_text, model_name, is_diagram=False, diagram_text=""):
    system_prompt = """
    You are an expert in ADA compliance for industrial and systems engineering courses. 
    Generate concise, pedagogical Alt Text (under 125 chars). 
    Focus on the logistics systems, material flow, batch processing, or facilities design depicted. 
    If relevant, describe sustainability or environmental themes in the context of the Engineering for One Planet framework.
    CRITICAL: If the image is an empty placeholder, a blank frame, or a generic PowerPoint 'click to add picture' icon, output EXACTLY: DECORATIVE.
    """
    
    if is_diagram:
        user_prompt = f"Describe this structural diagram based on its extracted text: '{diagram_text}'. Context from the rest of the slide: {curr_text}"
        contents = [user_prompt]
    else:
        image = Image.open(io.BytesIO(image_bytes))
        user_prompt = f"Analyze this image. Context: {curr_text}. Previous context: {prev_text}."
        contents = [image, user_prompt]

    if "2.5-flash" in model_name:
        target_rpm = 4.0
    else:
        target_rpm = 14.0 
        
    dynamic_sleep_time = 60.0 / target_rpm
    max_retries = 3 

    for attempt in range(max_retries):
        try:
            elapsed_time = time.time() - st.session_state.last_api_call
            if elapsed_time < dynamic_sleep_time:
                time.sleep(dynamic_sleep_time - elapsed_time)
            
            config_args = {
                "system_instruction": system_prompt,
                "temperature": 0.2
            }
            if "gemma-4" in model_name:
                config_args["thinking_config"] = types.ThinkingConfig(thinking_level="high")

            response = client.models.generate_content(
                model=model_name,
                contents=contents,
                config=types.GenerateContentConfig(**config_args)
            )
            
            st.session_state.last_api_call = time.time()
            return response.text.strip()
            
        except Exception as e:
            st.session_state.last_api_call = time.time() 
            err_str = str(e).lower()
            if "429" in err_str or "503" in err_str or "quota" in err_str:
                if "day" in err_str or "daily" in err_str:
                    return f"Error: Daily API quota exceeded for {model_name}."
                
                if attempt < max_retries - 1:
                    sec_match = re.search(r'in\s*(\d+)\s*s', err_str)
                    if not sec_match:
                        sec_match = re.search(r'(\d+)\s*second', err_str)
                        
                    wait_time = int(sec_match.group(1)) + 2 if sec_match else int(dynamic_sleep_time * (attempt + 2))
                    
                    countdown_placeholder = st.empty()
                    for seconds_left in range(wait_time, 0, -1):
                        countdown_placeholder.warning(f"⏳ Rate limit hit! Resuming in {seconds_left} seconds...")
                        time.sleep(1)
                    countdown_placeholder.empty() 
                    continue 
            return f"Error: {str(e)}"
    return "Error: Model overloaded."

def generate_and_add_title(client, slide, slide_text, model_name):
    has_title = False
    
    for shape in slide.shapes:
        if shape.is_placeholder and hasattr(shape, "placeholder_format"):
            if shape.placeholder_format.type in [1, 3]:
                if getattr(shape, "has_text_frame", False) and shape.text.strip():
                    has_title = True
                break 

    if not has_title:
        title_text = "Slide Content" 
        
        if slide_text.strip():
            prompt = f"Create a concise, 3-to-6 word title for a presentation slide containing this text. Output ONLY the title.\n\nText: {slide_text}"
            
            config_args = {"temperature": 0.2}
            if "gemma-4" in model_name:
                config_args["thinking_config"] = types.ThinkingConfig(thinking_level="high")
                
            for attempt in range(3):
                try:
                    time.sleep(2) 
                    response = client.models.generate_content(
                        model=model_name, 
                        contents=prompt,
                        config=types.GenerateContentConfig(**config_args)
                    )
                    title_text = response.text.strip()
                    break
                except Exception:
                    time.sleep(4)
        else:
            title_text = "Visual Presentation Slide"

        try:
            txBox = slide.shapes.add_textbox(Inches(-5), Inches(-5), Inches(1), Inches(1))
            txBox.text = title_text
            force_textbox_to_title(txBox) 
        except Exception:
            pass

# --- Main App ---
st.set_page_config(page_title="ADA PPTX Automator Pro", layout="centered")
st.title("♿ ADA Course Material Automator")
st.markdown("Features deep XML native injection, Safe Ghost Overlays for PDF compliance, and Context-Aware AI Title generation.")

api_key = st.text_input("Enter your Gemini API Key:", type="password")

st.markdown("### Model Selection")
selected_model = st.selectbox(
    "Choose AI Model:",
    ("gemma-4-31b-it", "gemini-2.5-flash", "gemini-1.5-flash")
)

st.markdown("### Select ADA Fixes to Apply:")
do_captions = st.checkbox("Generate Image/SmartArt Captions (Alt Text)", value=True)
do_titles = st.checkbox("Generate Missing Slide Titles", value=True)
do_reading_order = st.checkbox("Fix Reading Order (Top-to-Bottom)", value=True)

uploaded_file = st.file_uploader("Upload PowerPoint File", type=["pptx"])

if uploaded_file and api_key:
    if st.button("Process Presentation"):
        client = genai.Client(api_key=api_key)
        prs = Presentation(uploaded_file)
        
        saved_calls = 0
        api_calls = 0
        ghost_shapes = 0
        prev_text = ""
        
        with st.spinner(f"Processing slides using {selected_model}..."):
            progress_bar = st.progress(0)
            total_slides = len(prs.slides)
            
            for i, slide in enumerate(prs.slides):
                curr_text = get_slide_text(slide) 
                
                if do_titles:
                    generate_and_add_title(client, slide, curr_text, selected_model)
                
                if do_captions:
                    for shape in slide.shapes:
                        if shape.is_placeholder:
                            shape_has_text = bool(get_shape_text(shape).strip())
                            shape_has_image = False
                            try:
                                if hasattr(shape, "image") and shape.image.blob:
                                    shape_has_image = True
                            except Exception:
                                pass
                                
                            if not shape_has_text and not shape_has_image:
                                mark_as_decorative(shape) 
                                ghost_shapes += 1
                                continue
                                
                        try:
                            if hasattr(shape, "image"):
                                img_bytes = shape.image.blob
                                img_hash = shape.image.sha1 
                                
                                if img_hash in st.session_state.caption_cache:
                                    set_alt_text(shape, st.session_state.caption_cache[img_hash])
                                    saved_calls += 1
                                else:
                                    caption = generate_caption(client, img_bytes, prev_text, curr_text, model_name=selected_model)
                                    if not caption.startswith("Error"):
                                        st.session_state.caption_cache[img_hash] = caption
                                        if caption.upper() == "DECORATIVE":
                                            mark_as_decorative(shape)
                                        else:
                                            set_alt_text(shape, caption)
                                        api_calls += 1
                                    else:
                                        st.warning(f"Slide {i+1} Issue: {caption}")
                                continue 
                        except Exception:
                            pass
                        
                        if getattr(shape, "shape_type", None) in [MSO_SHAPE_TYPE.GROUP, 19, 24] or hasattr(shape._element, 'nvGraphicFramePr'):
                            d_text = get_shape_text(shape)
                            caption = generate_caption(client, None, prev_text, curr_text, model_name=selected_model, is_diagram=True, diagram_text=d_text)
                            
                            if not caption.startswith("Error"):
                                mute_smartart_children(shape) 
                                mark_as_decorative(shape)
                                create_ghost_overlay(slide, shape, caption) 
                                api_calls += 1
                            else:
                                st.warning(f"Slide {i+1} SmartArt Issue: {caption}")
                                    
                if do_reading_order:
                    fix_reading_order(slide)
                
                prev_text = curr_text
                progress_bar.progress((i + 1) / total_slides)

            output = io.BytesIO()
            prs.save(output)
            output.seek(0)
            
            st.success(f"Finished! API Calls: {api_calls} | Redundant Images Saved: {saved_calls} | Ghost Shapes Muted: {ghost_shapes}")
            st.download_button("Download ADA File", output, file_name=f"ADA_Compliant_{uploaded_file.name}")